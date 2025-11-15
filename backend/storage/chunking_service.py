"""
Document chunking service for RAG implementation.
Splits documents into optimal chunks for embedding and retrieval.

Supports multiple chunking strategies inspired by Gemini File Search:
- Auto: Automatically choose best strategy based on content
- Whitespace: Split on whitespace boundaries
- Semantic: Split on semantic boundaries (paragraphs, sections)
- Fixed: Fixed-size chunks with overlap
"""

import logging
import os
import re
from typing import List, Dict, Tuple, Optional
from pathlib import Path

logger = logging.getLogger(__name__)


class ChunkingService:
    """
    Service for chunking documents into optimal sizes for embedding.
    Implements Gemini-style configurable chunking strategies.
    """

    # Approximate characters per token (rough estimate)
    CHARS_PER_TOKEN = 4

    def __init__(
        self,
        max_chunk_size: int = 500,
        overlap_size: int = 50,
        separators: List[str] = None,
        strategy: str = 'auto',
        max_tokens_per_chunk: int = 512,
        max_overlap_tokens: int = 50
    ):
        """
        Initialize chunking service.

        Args:
            max_chunk_size: Maximum number of characters per chunk
            overlap_size: Number of overlapping characters between chunks
            separators: List of separators to split on (in order of preference)
            strategy: Chunking strategy ('auto', 'whitespace', 'semantic', 'fixed')
            max_tokens_per_chunk: Maximum tokens per chunk (for token-based chunking)
            max_overlap_tokens: Overlap in tokens between chunks
        """
        self.max_chunk_size = max_chunk_size
        self.overlap_size = overlap_size
        self.separators = separators or ['\n\n', '\n', '. ', ' ', '']
        self.strategy = strategy
        self.max_tokens_per_chunk = max_tokens_per_chunk
        self.max_overlap_tokens = max_overlap_tokens

    def chunk_text(self, text: str, metadata: Dict = None, strategy: Optional[str] = None) -> List[Dict]:
        """
        Split text into chunks with overlap using specified strategy.

        Args:
            text: Input text to chunk
            metadata: Optional metadata to attach to each chunk
            strategy: Chunking strategy override ('auto', 'whitespace', 'semantic', 'fixed')

        Returns:
            List of chunk dictionaries with text and metadata
        """
        if not text or not text.strip():
            return []

        chunks = []
        metadata = metadata or {}
        strategy = strategy or self.strategy

        # Choose chunking strategy
        if strategy == 'auto':
            current_chunks = self._auto_chunk(text)
        elif strategy == 'whitespace':
            current_chunks = self._whitespace_chunk(text)
        elif strategy == 'semantic':
            current_chunks = self._semantic_chunk(text)
        elif strategy == 'fixed':
            current_chunks = self._fixed_chunk(text)
        else:
            # Fallback to recursive split
            current_chunks = self._recursive_split(text, self.separators)

        # Process chunks with overlap
        for i, chunk_text in enumerate(current_chunks):
            if not chunk_text.strip():
                continue

            # Estimate token count
            token_count = self._estimate_tokens(chunk_text)

            chunk_dict = {
                'chunk_index': i,
                'chunk_text': chunk_text.strip(),
                'chunk_size': len(chunk_text),
                'token_count': token_count,
                'chunking_strategy': strategy,
                'overlap_tokens': self.max_overlap_tokens if i > 0 else 0,
                'metadata': metadata.copy()
            }
            chunks.append(chunk_dict)

        return chunks

    def _estimate_tokens(self, text: str) -> int:
        """
        Estimate token count for text.

        Args:
            text: Text to estimate

        Returns:
            Estimated token count
        """
        # Simple estimation: ~4 characters per token
        return len(text) // self.CHARS_PER_TOKEN

    def _auto_chunk(self, text: str) -> List[str]:
        """
        Automatically choose best chunking strategy based on content.

        Args:
            text: Text to chunk

        Returns:
            List of text chunks
        """
        # Analyze text structure
        has_paragraphs = '\n\n' in text
        has_sections = bool(re.search(r'^#{1,6}\s+', text, re.MULTILINE))  # Markdown headers
        avg_line_length = len(text) / (text.count('\n') + 1) if text.count('\n') > 0 else len(text)

        # Choose strategy
        if has_sections or has_paragraphs:
            return self._semantic_chunk(text)
        elif avg_line_length > 100:
            return self._semantic_chunk(text)
        else:
            return self._whitespace_chunk(text)

    def _whitespace_chunk(self, text: str) -> List[str]:
        """
        Chunk text on whitespace boundaries with token limits.

        Args:
            text: Text to chunk

        Returns:
            List of text chunks
        """
        max_chars = self.max_tokens_per_chunk * self.CHARS_PER_TOKEN
        overlap_chars = self.max_overlap_tokens * self.CHARS_PER_TOKEN

        chunks = []
        words = text.split()
        current_chunk = []
        current_length = 0

        for word in words:
            word_length = len(word) + 1  # +1 for space

            if current_length + word_length > max_chars and current_chunk:
                # Save current chunk
                chunks.append(' '.join(current_chunk))

                # Create overlap by keeping last few words
                overlap_words = []
                overlap_length = 0
                for w in reversed(current_chunk):
                    if overlap_length + len(w) + 1 <= overlap_chars:
                        overlap_words.insert(0, w)
                        overlap_length += len(w) + 1
                    else:
                        break

                current_chunk = overlap_words
                current_length = overlap_length

            current_chunk.append(word)
            current_length += word_length

        # Add final chunk
        if current_chunk:
            chunks.append(' '.join(current_chunk))

        return chunks

    def _semantic_chunk(self, text: str) -> List[str]:
        """
        Chunk text on semantic boundaries (paragraphs, sections).

        Args:
            text: Text to chunk

        Returns:
            List of text chunks
        """
        max_chars = self.max_tokens_per_chunk * self.CHARS_PER_TOKEN
        overlap_chars = self.max_overlap_tokens * self.CHARS_PER_TOKEN

        # Split on semantic boundaries
        # Priority: sections (###), paragraphs (\n\n), sentences, words
        semantic_separators = [
            r'\n#{1,6}\s+[^\n]+\n',  # Markdown headers
            '\n\n',                    # Paragraphs
            '. ',                      # Sentences
            '\n',                      # Lines
            ' '                        # Words
        ]

        return self._recursive_split_semantic(text, semantic_separators, max_chars, overlap_chars)

    def _recursive_split_semantic(self, text: str, separators: List[str], max_chars: int, overlap_chars: int) -> List[str]:
        """
        Recursively split text using semantic separators.

        Args:
            text: Text to split
            separators: List of regex patterns or strings to split on
            max_chars: Maximum characters per chunk
            overlap_chars: Overlap characters

        Returns:
            List of text chunks
        """
        if not separators or len(text) <= max_chars:
            return [text] if text else []

        separator = separators[0]
        remaining_separators = separators[1:]

        # Try to split by current separator
        if separator.startswith(r'\\n'):
            # Regex pattern
            splits = re.split(separator, text)
        else:
            splits = text.split(separator)

        chunks = []
        current_chunk = ""

        for i, split in enumerate(splits):
            if not split:
                continue

            # Reconstruct with separator
            piece = split if i == 0 or not current_chunk else separator + split

            if len(current_chunk) + len(piece) <= max_chars:
                current_chunk += piece
            else:
                # Current chunk is full
                if current_chunk:
                    chunks.append(current_chunk)

                # If split is too large, recursively split
                if len(piece) > max_chars:
                    sub_chunks = self._recursive_split_semantic(piece, remaining_separators, max_chars, overlap_chars)
                    chunks.extend(sub_chunks)
                    current_chunk = ""
                else:
                    current_chunk = piece

        if current_chunk:
            chunks.append(current_chunk)

        return chunks

    def _fixed_chunk(self, text: str) -> List[str]:
        """
        Create fixed-size chunks with token-based overlap.

        Args:
            text: Text to chunk

        Returns:
            List of text chunks
        """
        max_chars = self.max_tokens_per_chunk * self.CHARS_PER_TOKEN
        overlap_chars = self.max_overlap_tokens * self.CHARS_PER_TOKEN

        chunks = []
        start = 0

        while start < len(text):
            end = start + max_chars
            chunk = text[start:end]

            # Try to break on word boundary
            if end < len(text):
                # Find last space before end
                last_space = chunk.rfind(' ')
                if last_space > max_chars * 0.8:  # Only if we're not losing too much
                    chunk = chunk[:last_space]
                    end = start + last_space

            chunks.append(chunk)

            # Move start forward with overlap
            start = end - overlap_chars

        return chunks

    def _recursive_split(self, text: str, separators: List[str]) -> List[str]:
        """
        Recursively split text using separators.

        Args:
            text: Text to split
            separators: List of separators to try

        Returns:
            List of text chunks
        """
        if not separators:
            # Base case: split by character limit
            return self._split_by_length(text)

        separator = separators[0]
        remaining_separators = separators[1:]

        if separator == '':
            # Last resort: split by character
            return self._split_by_length(text)

        splits = text.split(separator)
        chunks = []

        current_chunk = ""
        for split in splits:
            # Try to add split to current chunk
            test_chunk = current_chunk + separator + split if current_chunk else split

            if len(test_chunk) <= self.max_chunk_size:
                current_chunk = test_chunk
            else:
                # Current chunk is full
                if current_chunk:
                    chunks.append(current_chunk)

                # If the split itself is too large, recursively split it
                if len(split) > self.max_chunk_size:
                    sub_chunks = self._recursive_split(split, remaining_separators)
                    chunks.extend(sub_chunks)
                    current_chunk = ""
                else:
                    current_chunk = split

        # Add remaining chunk
        if current_chunk:
            chunks.append(current_chunk)

        return chunks

    def _split_by_length(self, text: str) -> List[str]:
        """
        Split text by fixed length with overlap.

        Args:
            text: Text to split

        Returns:
            List of text chunks
        """
        chunks = []
        start = 0

        while start < len(text):
            end = start + self.max_chunk_size
            chunk = text[start:end]
            chunks.append(chunk)

            # Move start forward with overlap
            start = end - self.overlap_size

        return chunks

    def extract_text_from_file(self, file_path: str, file_type: str = None) -> str:
        """
        Extract text content from various file types.
        Supports Gemini-style wide file format support.

        Args:
            file_path: Path to the file
            file_type: Optional file type hint

        Returns:
            Extracted text content
        """
        file_ext = Path(file_path).suffix.lower()

        try:
            # Plain text files
            if file_ext in ['.txt', '.md', '.csv', '.log', '.json', '.xml', '.yaml', '.yml', '.rtf']:
                return self._read_text_file(file_path)

            # Code files (extended support)
            elif file_ext in ['.py', '.js', '.java', '.cpp', '.c', '.h', '.cs', '.php',
                            '.rb', '.go', '.rs', '.swift', '.kt', '.ts', '.jsx', '.tsx',
                            '.dart', '.sql', '.sh', '.bash', '.zsh', '.ps1', '.r', '.m',
                            '.scala', '.clj', '.ex', '.exs', '.erl', '.hrl', '.lua',
                            '.perl', '.pl', '.vb', '.asm', '.f90', '.jl']:
                return self._read_text_file(file_path)

            # PDF files
            elif file_ext == '.pdf':
                return self._extract_pdf_text(file_path)

            # Word documents
            elif file_ext in ['.docx', '.doc']:
                return self._extract_docx_text(file_path)

            # Excel spreadsheets
            elif file_ext in ['.xlsx', '.xls']:
                return self._extract_xlsx_text(file_path)

            # PowerPoint presentations
            elif file_ext in ['.pptx', '.ppt']:
                return self._extract_pptx_text(file_path)

            # HTML files
            elif file_ext in ['.html', '.htm']:
                return self._extract_html_text(file_path)

            # LaTeX files
            elif file_ext in ['.tex', '.latex']:
                return self._read_text_file(file_path)

            else:
                logger.warning(f"Unsupported file type for text extraction: {file_ext}")
                return ""

        except Exception as e:
            logger.error(f"Failed to extract text from {file_path}: {str(e)}")
            return ""

    def _read_text_file(self, file_path: str) -> str:
        """Read plain text file."""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

    def _extract_pdf_text(self, file_path: str) -> str:
        """Extract text from PDF file."""
        try:
            import PyPDF2
            text = ""
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n\n"
            return text
        except ImportError:
            logger.warning("PyPDF2 not installed. Cannot extract PDF text.")
            return ""
        except Exception as e:
            logger.error(f"PDF extraction failed: {str(e)}")
            return ""

    def _extract_docx_text(self, file_path: str) -> str:
        """Extract text from DOCX file."""
        try:
            import docx
            doc = docx.Document(file_path)
            text = "\n\n".join([paragraph.text for paragraph in doc.paragraphs])
            return text
        except ImportError:
            logger.warning("python-docx not installed. Cannot extract DOCX text.")
            return ""
        except Exception as e:
            logger.error(f"DOCX extraction failed: {str(e)}")
            return ""

    def _extract_html_text(self, file_path: str) -> str:
        """Extract text from HTML file."""
        try:
            from bs4 import BeautifulSoup
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
                # Remove script and style elements
                for script in soup(["script", "style"]):
                    script.decompose()
                return soup.get_text(separator='\n', strip=True)
        except ImportError:
            logger.warning("BeautifulSoup not installed. Cannot extract HTML text.")
            return self._read_text_file(file_path)  # Fallback to plain text
        except Exception as e:
            logger.error(f"HTML extraction failed: {str(e)}")
            return ""

    def _extract_xlsx_text(self, file_path: str) -> str:
        """
        Extract text from Excel (.xlsx, .xls) files.

        Args:
            file_path: Path to Excel file

        Returns:
            Extracted text from all sheets
        """
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            text_parts = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"\n=== Sheet: {sheet_name} ===\n")

                for row in sheet.iter_rows(values_only=True):
                    # Filter out None values and convert to strings
                    row_values = [str(cell) for cell in row if cell is not None]
                    if row_values:
                        text_parts.append(" | ".join(row_values))

            return "\n".join(text_parts)

        except ImportError:
            logger.warning("openpyxl not installed. Cannot extract XLSX text.")
            return ""
        except Exception as e:
            logger.error(f"XLSX extraction failed: {str(e)}")
            return ""

    def _extract_pptx_text(self, file_path: str) -> str:
        """
        Extract text from PowerPoint (.pptx, .ppt) files.

        Args:
            file_path: Path to PowerPoint file

        Returns:
            Extracted text from all slides
        """
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_parts = []

            for i, slide in enumerate(prs.slides, 1):
                text_parts.append(f"\n=== Slide {i} ===\n")

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text)

                    # Extract text from tables
                    if hasattr(shape, "table"):
                        for row in shape.table.rows:
                            row_text = " | ".join(cell.text for cell in row.cells)
                            if row_text.strip():
                                text_parts.append(row_text)

            return "\n".join(text_parts)

        except ImportError:
            logger.warning("python-pptx not installed. Cannot extract PPTX text.")
            return ""
        except Exception as e:
            logger.error(f"PPTX extraction failed: {str(e)}")
            return ""


# Singleton instance
chunking_service = ChunkingService()
